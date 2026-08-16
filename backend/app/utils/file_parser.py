"""
文档解析器 - 支持 PDF, DOCX, PPTX, Markdown, HTML, EPUB, TXT, 图像 (v3.2 Q3)
修复：PyPDF2 → pypdf、编码检测、流式读取、容错解析
v3.2: + ImageParser (PaddleOCR + BLIP Captioning)
"""
import logging
import os
from pathlib import Path

logger = logging.getLogger(__name__)

# 支持的文件扩展名 → MIME 类型映射 (v3.2: + 图像类型)
EXT_TO_TYPE: dict = {
    ".txt": "Text",
    ".md": "Markdown",
    ".markdown": "Markdown",
    ".pdf": "PDF",
    ".docx": "Word",
    ".pptx": "PowerPoint",
    ".html": "HTML",
    ".htm": "HTML",
    ".epub": "EPUB",
    ".jpg": "Image",
    ".jpeg": "Image",
    ".png": "Image",
    ".webp": "Image",
}

# 图像扩展名集合
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp"}


def read_file_content(filepath: str) -> str | None:
    """
    解析文件内容，根据扩展名调用对应的解析器。
    返回文件文本内容，解析失败返回 None。
    """
    ext = Path(filepath).suffix.lower()

    if not os.path.exists(filepath):
        logger.error(f"文件不存在: {filepath}")
        return None

    if os.path.getsize(filepath) == 0:
        logger.warning(f"文件为空: {filepath}")
        return None

    parsers = {
        ".txt": _read_text_with_encoding_detect,
        ".md": _read_text_with_encoding_detect,
        ".markdown": _read_text_with_encoding_detect,
        ".html": _read_text_with_encoding_detect,
        ".htm": _read_text_with_encoding_detect,
        ".pdf": _read_pdf,
        ".docx": _read_docx,
        ".pptx": _read_pptx,
        ".epub": _read_epub,
        ".jpg": _read_image,
        ".jpeg": _read_image,
        ".png": _read_image,
        ".webp": _read_image,
    }

    parser = parsers.get(ext)
    if parser is None:
        logger.warning(f"不支持的文件类型: {ext}")
        return None

    try:
        return parser(filepath)
    except Exception as e:
        logger.error(f"解析文件失败 {filepath}: {e}", exc_info=True)
        return None


# ─── 解压炸弹/超大输入防护 (v4.1) ──────────────────────────────────


def _check_zip_safety(filepath: str) -> None:
    """
    校验 zip 类容器（docx/pptx/epub）的安全性：
    - 总解压大小上限（PARSE_MAX_ZIP_TOTAL_MB）
    - 单文件压缩比上限（PARSE_MAX_ZIP_RATIO），拦截解压炸弹

    Raises:
        ValueError: 超出安全限制时抛出，由 read_file_content 捕获后拒绝解析
    """
    import zipfile

    from app.core.config import config
    try:
        with zipfile.ZipFile(filepath) as zf:
            total = sum(info.file_size for info in zf.infolist())
            limit_bytes = config.PARSE_MAX_ZIP_TOTAL_MB * 1024 * 1024
            if total > limit_bytes:
                raise ValueError(
                    f"压缩包总解压大小 ({total / 1024 / 1024:.1f}MB) 超过限制 "
                    f"({config.PARSE_MAX_ZIP_TOTAL_MB}MB)，已拒绝解析"
                )
            for info in zf.infolist():
                if info.compress_size > 0 and info.file_size / info.compress_size > config.PARSE_MAX_ZIP_RATIO:
                    raise ValueError(
                        f"压缩比超限 ({info.file_size}/{info.compress_size} > "
                        f"{config.PARSE_MAX_ZIP_RATIO})，疑似解压炸弹，已拒绝解析"
                    )
    except zipfile.BadZipFile:
        # 非 zip 结构交给上层解析器按原逻辑报错
        pass


def _enforce_total_chars_limit(texts: list, filepath: str) -> None:
    """累积文本量超过 PARSE_MAX_TEXT_CHARS 时截断并告警（就地修改 texts）"""
    from app.core.config import config

    total = sum(len(t) for t in texts)
    if total <= config.PARSE_MAX_TEXT_CHARS:
        return
    kept, acc = [], 0
    for t in texts:
        if acc >= config.PARSE_MAX_TEXT_CHARS:
            break
        kept.append(t)
        acc += len(t)
    texts[:] = kept
    logger.warning(
        f"解析文本量超过 {config.PARSE_MAX_TEXT_CHARS} 字符上限，已截断: {filepath}"
    )


# ─── 纯文本（含编码自动检测）───────────────────────────────────────


def _read_text_with_encoding_detect(filepath: str) -> str:
    """
    读取纯文本文件，自动检测编码。
    优先 UTF-8 → GBK → GB2312 → GB18030 → latin-1。
    v4.1: 读取字节量设上限，防止超大文本撑爆内存。
    """
    from app.core.config import config

    max_bytes = config.PARSE_MAX_TEXT_CHARS * 4
    with open(filepath, "rb") as f:
        raw = f.read(max_bytes + 1)
    if len(raw) > max_bytes:
        logger.warning(f"文本文件超过解析上限，已截断: {filepath}")
        raw = raw[:max_bytes]

    # 尝试 UTF-8（截断可能切断多字节字符，用 ignore 容错）
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        pass

    # 尝试常见中文编码
    for enc in ["gbk", "gb2312", "gb18030"]:
        try:
            text = raw.decode(enc)
            logger.info(f"检测到编码 {enc}: {filepath}")
            return text
        except UnicodeDecodeError:
            continue

    # 最终回退
    logger.warning(f"使用 latin-1 回退编码: {filepath}")
    return raw.decode("latin-1", errors="replace")


# ─── PDF（使用 pypdf）───────────────────────────────────────────────


def _read_pdf(filepath: str) -> str:
    """
    读取 PDF 文件 (v2.4: 仅使用 pypdf, 移除已废弃的 PyPDF2 回退)。
    v4.1: 页数与总文本量设上限。
    """
    from pypdf import PdfReader

    from app.core.config import config

    reader = PdfReader(filepath)
    texts = []
    failed_pages = 0

    all_pages = reader.pages
    if len(all_pages) > config.PARSE_MAX_PDF_PAGES:
        logger.warning(
            f"PDF 页数 ({len(all_pages)}) 超过上限 {config.PARSE_MAX_PDF_PAGES}，仅解析前 {config.PARSE_MAX_PDF_PAGES} 页: {filepath}"
        )
        all_pages = all_pages[: config.PARSE_MAX_PDF_PAGES]

    total_chars = 0
    for i, page in enumerate(all_pages):
        try:
            text = page.extract_text()
            if text:
                stripped = text.strip()
                texts.append(stripped)
                total_chars += len(stripped)
                if total_chars > config.PARSE_MAX_TEXT_CHARS:
                    logger.warning(f"PDF 文本量超过 {config.PARSE_MAX_TEXT_CHARS} 字符上限，已截断: {filepath}")
                    break
        except Exception as e:
            failed_pages += 1
            logger.debug(f"PDF 第 {i} 页解析失败: {e}")

    if failed_pages:
        logger.warning(f"PDF {filepath}: {failed_pages}/{len(all_pages)} 页解析失败")

    if not texts:
        logger.error(f"PDF 无有效文本: {filepath}")
        return ""

    return "\n\n".join(texts)


# ─── DOCX ─────────────────────────────────────────────────────────


def _read_docx(filepath: str) -> str:
    """读取 Word 文档，保留标题层级（v4.1: 解压炸弹防护）"""
    _check_zip_safety(filepath)
    from docx import Document

    doc = Document(filepath)
    texts = []
    for para in doc.paragraphs:
        line = para.text.strip()
        if not line:
            continue
        # 根据样式区分标题和正文
        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            level_match = style_name.split()[-1]
            try:
                level_num = int(level_match)
                prefix = "#" * min(level_num, 6)
            except ValueError:
                prefix = "#"
            texts.append(f"{prefix} {line}")
        else:
            texts.append(line)
    return "\n\n".join(texts)


# ─── PPTX ─────────────────────────────────────────────────────────


def _read_pptx(filepath: str) -> str:
    """读取 PowerPoint 文件，提取所有幻灯片文本（v4.1: 解压炸弹防护）"""
    _check_zip_safety(filepath)
    from pptx import Presentation

    prs = Presentation(filepath)
    texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"## 幻灯片 {slide_num}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)

            # 提取表格内容
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))

        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


# ─── EPUB ─────────────────────────────────────────────────────────


def _read_epub(filepath: str) -> str:
    """读取 EPUB 电子书（v4.1: 解压炸弹防护 + 总文本量上限）"""
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError:
        logger.warning("ebooklib 未安装，无法解析 EPUB")
        return ""

    _check_zip_safety(filepath)
    book = epub.read_epub(filepath)
    texts = []

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            try:
                content = item.get_content().decode("utf-8")
                # 简单的 HTML 标签剥离
                clean = _strip_html(content)
                if clean.strip():
                    texts.append(clean.strip())
            except Exception as e:
                logger.debug(f"EPUB 章节解析失败: {e}")

    _enforce_total_chars_limit(texts, filepath)
    return "\n\n".join(texts)


# ─── HTML ─────────────────────────────────────────────────────────


def _strip_html(html: str) -> str:
    """简易 HTML 标签剥离"""
    import re
    # 移除 script 和 style
    html = re.sub(r'<(script|style)[^>]*>.*?</\1>', '', html, flags=re.DOTALL | re.IGNORECASE)
    # 移除 HTML 标签
    html = re.sub(r'<[^>]+>', ' ', html)
    # 解码常见实体
    html = html.replace("&nbsp;", " ").replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    html = html.replace("&quot;", '"').replace("&#39;", "'")
    # 压缩空白
    html = re.sub(r'[ \t]+', ' ', html)
    html = re.sub(r'\n{3,}', '\n\n', html)
    return html.strip()


# ─── 图像 (PaddleOCR + BLIP Captioning) ────────────────────────────

# 模块级模型缓存（懒加载）
_ocr_model = None
_blip_model = None
_blip_processor = None


def _get_ocr():
    """懒加载 PaddleOCR 模型"""
    global _ocr_model
    if _ocr_model is None:
        try:
            from paddleocr import PaddleOCR
            _ocr_model = PaddleOCR(lang='ch', use_angle_cls=False, show_log=False)
            logger.info("PaddleOCR 模型加载完成")
        except ImportError:
            logger.warning("PaddleOCR 未安装，图像 OCR 功能不可用")
            return None
    return _ocr_model


def _get_blip():
    """懒加载 BLIP 图像描述模型"""
    global _blip_model, _blip_processor
    if _blip_model is None:
        try:
            from transformers import BlipForConditionalGeneration, BlipProcessor
            _blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
            _blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
            logger.info("BLIP 图像描述模型加载完成")
        except ImportError:
            logger.warning("transformers 未安装，BLIP 图像描述功能不可用")
            return None, None
        except Exception as e:
            logger.warning(f"BLIP 模型加载失败: {e}")
            return None, None
    return _blip_processor, _blip_model


def _read_image(filepath: str) -> str:
    """
    读取图像文件，提取文本和语义信息。

    使用 PaddleOCR 提取文字 + BLIP 生成图片描述。
    返回拼接的文本: [OCR文字]：{ocr_text}。 [图片描述]：{caption_text}
    """
    from PIL import Image

    parts = []

    # Step 1: OCR 文字提取
    ocr = _get_ocr()
    if ocr is not None:
        try:
            result = ocr.ocr(filepath)
            ocr_texts = []
            if result and result[0]:
                for line in result[0]:
                    if line and len(line) >= 2:
                        text = line[1][0] if isinstance(line[1], (list, tuple)) else str(line[1])
                        if text and text.strip():
                            ocr_texts.append(text.strip())
            if ocr_texts:
                parts.append(f"[OCR文字]：{'；'.join(ocr_texts)}")
        except Exception as e:
            logger.warning(f"OCR 识别失败 {filepath}: {e}")

    # Step 2: BLIP 图像描述
    processor, model = _get_blip()
    if processor is not None and model is not None:
        try:
            image = Image.open(filepath).convert("RGB")
            inputs = processor(image, return_tensors="pt")
            out = model.generate(**inputs, max_new_tokens=50)
            caption = processor.decode(out[0], skip_special_tokens=True)
            if caption and caption.strip():
                parts.append(f"[图片描述]：{caption.strip()}")
        except Exception as e:
            logger.warning(f"BLIP 图像描述失败 {filepath}: {e}")

    if not parts:
        logger.warning(f"图像无有效内容: {filepath}")
        return ""

    result = "。 ".join(parts)
    logger.info(f"图像解析完成: {filepath}, 内容长度={len(result)}")
    return result


# ─── 文件信息 ─────────────────────────────────────────────────────


def get_file_info(filepath: str, file_hash: str = None) -> dict:
    """获取文件基本信息（含 MIME 检测）。

    Args:
        filepath: 文件路径
        file_hash: 预计算的 SHA256 哈希（避免重复计算）
    """
    path = Path(filepath)
    ext = path.suffix.lower()

    from app.utils.helpers import compute_file_hash, detect_mime_type

    size = path.stat().st_size if path.exists() else 0
    mime = detect_mime_type(filepath)
    if file_hash is None:
        file_hash = compute_file_hash(filepath)

    return {
        "name": path.name,
        "path": str(path),
        "size": size,
        "type": EXT_TO_TYPE.get(ext, "Unknown"),
        "extension": ext,
        "mime_type": mime,
        "sha256": file_hash,
    }
